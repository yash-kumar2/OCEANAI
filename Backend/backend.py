import os
import json
import io
import datetime
import re  # Added regex for parsing
from bson.objectid import ObjectId
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from flask_jwt_extended import JWTManager, create_access_token, jwt_required, get_jwt_identity
from flask_bcrypt import Bcrypt
from pymongo import MongoClient
import google.generativeai as genai
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Configuration ---
app = Flask(__name__)
CORS(app)

# ---------------------------------------------------------
# [CONFIGURATION AREA]
# ---------------------------------------------------------
MONGO_URI = os.environ.get("MONGO_URI", "")
app.config["JWT_SECRET_KEY"] = os.environ.get("JWT_SECRET_KEY", "super-secret-change-me-in-prod")
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")

# ---------------------------------------------------------

# Service Initialization
bcrypt = Bcrypt(app)
jwt = JWTManager(app)
genai.configure(api_key=GEMINI_API_KEY)

# Database Setup
if not MONGO_URI:
    print("WARNING: MONGO_URI is not set. Database operations will fail.")
    client = None
    db = None
else:
    try:
        client = MongoClient(MONGO_URI)
        db = client.get_database("ai_authoring_platform")
        print("Connected to MongoDB.")
    except Exception as e:
        print(f"Failed to connect to MongoDB: {e}")
        client = None
        db = None

# Model Config
generation_config = {
    "temperature": 0.7,
    "top_p": 1,
    "top_k": 1,
    "max_output_tokens": 2048,
}
model = genai.GenerativeModel("gemini-2.5-flash-preview-09-2025", generation_config=generation_config)

# --- Helper Functions ---

def clean_json_text(text):
    text = text.strip()
    if text.startswith("```json"):
        text = text[7:]
    if text.endswith("```"):
        text = text[:-3]
    return text.strip()

def serialize_doc(doc):
    if not doc: return None
    doc["_id"] = str(doc["_id"])
    return doc

def add_markdown_text_to_paragraph(paragraph, text):
    """
    Parses text with **bold** markers and adds runs to a python-docx or python-pptx paragraph.
    """
    # Split by ** markers
    parts = text.split('**')
    for i, part in enumerate(parts):
        run = paragraph.add_run()
        run.text = part
        # Odd indices (1, 3, 5...) are the ones between ** markers
        if i % 2 == 1:
            run.font.bold = True

# --- Auth Routes (Unchanged) ---
@app.route('/api/register', methods=['POST'])
def register():
    data = request.json
    email = data.get('email')
    password = data.get('password')
    if not email or not password: return jsonify({"error": "Email and password required"}), 400
    if db.users.find_one({"email": email}): return jsonify({"error": "User already exists"}), 400
    hashed_password = bcrypt.generate_password_hash(password).decode('utf-8')
    user_id = db.users.insert_one({"email": email, "password": hashed_password, "created_at": datetime.datetime.utcnow()}).inserted_id
    access_token = create_access_token(identity=str(user_id))
    return jsonify({"token": access_token, "email": email}), 201

@app.route('/api/login', methods=['POST'])
def login():
    data = request.json
    email = data.get('email')
    password = data.get('password')
    user = db.users.find_one({"email": email})
    if user and bcrypt.check_password_hash(user['password'], password):
        access_token = create_access_token(identity=str(user['_id']))
        return jsonify({"token": access_token, "email": email}), 200
    return jsonify({"error": "Invalid credentials"}), 401

# --- Project Management Routes (Unchanged) ---
@app.route('/api/projects', methods=['GET'])
@jwt_required()
def get_projects():
    user_id = get_jwt_identity()
    projects = list(db.projects.find({"user_id": user_id}).sort("created_at", -1))
    return jsonify([serialize_doc(p) for p in projects])

@app.route('/api/projects', methods=['POST'])
@jwt_required()
def create_project():
    user_id = get_jwt_identity()
    data = request.json
    new_project = {
        "user_id": user_id,
        "topic": data.get('topic'),
        "type": data.get('type'),
        "sections": [],
        "created_at": datetime.datetime.utcnow(),
        "last_modified": datetime.datetime.utcnow()
    }
    result = db.projects.insert_one(new_project)
    new_project["_id"] = result.inserted_id
    return jsonify(serialize_doc(new_project)), 201

@app.route('/api/projects/<project_id>', methods=['GET'])
@jwt_required()
def get_project(project_id):
    user_id = get_jwt_identity()
    project = db.projects.find_one({"_id": ObjectId(project_id), "user_id": user_id})
    if not project: return jsonify({"error": "Project not found"}), 404
    return jsonify(serialize_doc(project))

@app.route('/api/projects/<project_id>', methods=['PUT'])
@jwt_required()
def update_project(project_id):
    user_id = get_jwt_identity()
    data = request.json
    update_data = {"last_modified": datetime.datetime.utcnow()}
    if 'sections' in data: update_data['sections'] = data['sections']
    result = db.projects.update_one({"_id": ObjectId(project_id), "user_id": user_id}, {"$set": update_data})
    if result.matched_count == 0: return jsonify({"error": "Project not found"}), 404
    return jsonify({"status": "updated"})

# --- AI Generation Routes (Updated for PPT Length) ---

@app.route('/api/projects/<project_id>/generate-outline', methods=['POST'])
@jwt_required()
def generate_outline(project_id):
    user_id = get_jwt_identity()
    project = db.projects.find_one({"_id": ObjectId(project_id), "user_id": user_id})
    if not project: return jsonify({"error": "Project not found"}), 404

    prompt = ""
    if project['type'] == 'pptx':
        prompt = f"""Create a presentation outline for: "{project['topic']}". 
        Generate 5 to 8 slide titles. Return ONLY raw JSON string array. Example: ["Title", "Intro"]"""
    else:
        prompt = f"""Create a report outline for: "{project['topic']}". 
        Generate 5 to 8 section headers. Return ONLY raw JSON string array. Example: ["Intro", "Body"]"""

    try:
        response = model.generate_content(prompt)
        cleaned = clean_json_text(response.text)
        titles = json.loads(cleaned)
        sections = [{"title": t, "content": "", "generated": False, "feedback": None, "comments": []} for t in titles]

        db.projects.update_one({"_id": ObjectId(project_id)}, {"$set": {"sections": sections, "last_modified": datetime.datetime.utcnow()}})
        return jsonify({"sections": sections})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/projects/<project_id>/sections/generate', methods=['POST'])
@jwt_required()
def generate_section(project_id):
    user_id = get_jwt_identity()
    data = request.json
    index = data.get('index')
    
    project = db.projects.find_one({"_id": ObjectId(project_id), "user_id": user_id})
    if not project: return jsonify({"error": "Project not found"}), 404

    try:
        section = project['sections'][index]
        context = "slide" if project['type'] == 'pptx' else "document section"
        
        prompt = f"""Write content for a {context} titled "{section['title']}" about "{project['topic']}". 
        Professional tone."""
        
        # STRICT LENGTH CONTROL FOR PPT
        if project['type'] == 'pptx': 
            prompt += " STRICT CONSTRAINT: Provide exactly 3 to 5 bullet points. Max 15 words per bullet. No intro or outro text. Use **bold** for key terms."
        else: 
            prompt += " Write 2-3 well-structured paragraphs. Use **bold** for emphasis."

        response = model.generate_content(prompt)
        content = response.text

        db.projects.update_one(
            {"_id": ObjectId(project_id)},
            {"$set": {f"sections.{index}.content": content, f"sections.{index}.generated": True, "last_modified": datetime.datetime.utcnow()}}
        )
        return jsonify({"content": content})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/projects/<project_id>/sections/refine', methods=['POST'])
@jwt_required()
def refine_section(project_id):
    user_id = get_jwt_identity()
    data = request.json
    index = data.get('index')
    instruction = data.get('instruction')
    
    project = db.projects.find_one({"_id": ObjectId(project_id), "user_id": user_id})
    if not project: return jsonify({"error": "Project not found"}), 404

    try:
        current_content = project['sections'][index]['content']
        prompt = f"""Original: "{current_content}"\nInstruction: {instruction}\nRewrite it. Return ONLY the new content. Keep strict length constraints if this is a slide."""
        
        response = model.generate_content(prompt)
        new_content = response.text

        db.projects.update_one(
            {"_id": ObjectId(project_id)},
            {"$set": {f"sections.{index}.content": new_content, "last_modified": datetime.datetime.utcnow()}}
        )
        return jsonify({"content": new_content})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --- Export (Updated for formatting) ---

@app.route('/api/projects/<project_id>/export', methods=['GET'])
@jwt_required()
def export_project(project_id):
    user_id = get_jwt_identity()
    project = db.projects.find_one({"_id": ObjectId(project_id), "user_id": user_id})
    if not project: return jsonify({"error": "Project not found"}), 404

    file_stream = io.BytesIO()
    topic = project['topic']
    sections = project['sections']

    if project['type'] == 'docx':
        document = Document()
        document.add_heading(topic, 0)
        for sec in sections:
            document.add_heading(sec['title'], level=1)
            p = document.add_paragraph()
            add_markdown_text_to_paragraph(p, sec['content']) # Handle bold

        document.save(file_stream)
        mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        ext = 'docx'

    elif project['type'] == 'pptx':
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic
        
        for sec in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sec['title']
            
            # Handle Text Frame
            tf = slide.placeholders[1].text_frame
            tf.clear() # Clear default placeholder text
            
            lines = sec['content'].split('\n')
            for line in lines:
                line = line.strip()
                if not line: continue
                # Remove markdown bullet chars if Gemini added them (e.g., "* ", "- ")
                if line.startswith('* ') or line.startswith('- '):
                    line = line[2:]
                
                p = tf.add_paragraph()
                p.level = 0
                add_markdown_text_to_paragraph(p, line) # Handle bold inside bullet

        prs.save(file_stream)
        mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        ext = 'pptx'

    file_stream.seek(0)
    return send_file(
        file_stream, 
        as_attachment=True, 
        download_name=f"{topic.replace(' ', '_')}.{ext}", 
        mimetype=mimetype
    )

if __name__ == '__main__':
    app.run(debug=True, port=5000)